from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_FILE = "Policy_Sarthi_Jury_Presentation.pptx"

COLORS = {
    "bg": RGBColor(248, 250, 252),
    "navy": RGBColor(24, 44, 77),
    "blue": RGBColor(21, 101, 192),
    "teal": RGBColor(0, 137, 123),
    "green": RGBColor(67, 160, 71),
    "orange": RGBColor(245, 124, 0),
    "white": RGBColor(255, 255, 255),
    "gray": RGBColor(89, 99, 110),
    "border": RGBColor(214, 222, 231),
    "soft_blue": RGBColor(232, 241, 252),
    "soft_green": RGBColor(234, 246, 238),
    "soft_orange": RGBColor(255, 244, 232),
}


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12.0), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = COLORS["navy"]

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.02), Inches(2.8), Inches(0.07))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["teal"]
    line.line.fill.background()

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12.1), Inches(0.5))
        p = sub.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.size = Pt(12)
        r.font.color.rgb = COLORS["gray"]


def add_footer(slide, text):
    foot = slide.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(12.0), Inches(0.25))
    p = foot.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.color.rgb = COLORS["gray"]


def add_card(slide, left, top, width, height, title, body, fill):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = COLORS["border"]

    t = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.10), width - Inches(0.3), Inches(0.35))
    p = t.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(14)
    r.font.color.rgb = COLORS["navy"]

    b = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.45), width - Inches(0.3), height - Inches(0.55))
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = body
    r.font.size = Pt(11)
    r.font.color.rgb = COLORS["navy"]


def add_bullets(slide, bullets, left=0.8, top=1.9, width=11.7, height=4.8, size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.bullet = True
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["navy"]


def add_flow_box(slide, x, y, w, h, label, fill):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = COLORS["border"]
    txt = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.14), Inches(w - 0.24), Inches(h - 0.28))
    p = txt.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = COLORS["navy"]


def connect(slide, x1, y1, x2, y2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = COLORS["teal"]
    c.line.width = Pt(2)
    c.line.end_arrowhead = True


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, COLORS["bg"])
    band = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["navy"]
    band.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12.0), Inches(1.4))
    p = t.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Policy Sarthi AI Agent"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = COLORS["navy"]
    p2 = t.text_frame.add_paragraph()
    p2.text = "Hospital Policy Intelligence + Ayushman Bharat Guidance"
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLORS["gray"]
    chip = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(4.35), Inches(5.2), Inches(0.75))
    chip.fill.solid()
    chip.fill.fore_color.rgb = COLORS["teal"]
    chip.line.fill.background()
    cp = s.shapes.add_textbox(Inches(1.05), Inches(4.57), Inches(4.7), Inches(0.3)).text_frame.paragraphs[0]
    cr = cp.add_run()
    cr.text = "Jury Deck | Product + Architecture + Scale Plan"
    cr.font.size = Pt(14)
    cr.font.bold = True
    cr.font.color.rgb = COLORS["white"]

    # Visual accents aligned to problem statement (document-heavy + multilingual care workflow)
    accent_left = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(3.95), Inches(2.2), Inches(1.15))
    accent_left.fill.solid()
    accent_left.fill.fore_color.rgb = COLORS["soft_blue"]
    accent_left.line.color.rgb = COLORS["border"]
    at1 = s.shapes.add_textbox(Inches(7.2), Inches(4.2), Inches(1.8), Inches(0.7)).text_frame
    ap1 = at1.paragraphs[0]
    ap1.alignment = PP_ALIGN.CENTER
    ar1 = ap1.add_run()
    ar1.text = "Policy Docs\nOCR + RAG"
    ar1.font.size = Pt(12)
    ar1.font.bold = True
    ar1.font.color.rgb = COLORS["navy"]

    accent_mid = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.45), Inches(3.95), Inches(1.7), Inches(1.15))
    accent_mid.fill.solid()
    accent_mid.fill.fore_color.rgb = COLORS["soft_green"]
    accent_mid.line.color.rgb = COLORS["border"]
    at2 = s.shapes.add_textbox(Inches(9.62), Inches(4.28), Inches(1.35), Inches(0.5)).text_frame
    ap2 = at2.paragraphs[0]
    ap2.alignment = PP_ALIGN.CENTER
    ar2 = ap2.add_run()
    ar2.text = "Voice +\nMultilingual"
    ar2.font.size = Pt(11)
    ar2.font.bold = True
    ar2.font.color.rgb = COLORS["navy"]

    accent_right = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.35), Inches(3.95), Inches(1.2), Inches(1.15))
    accent_right.fill.solid()
    accent_right.fill.fore_color.rgb = COLORS["soft_orange"]
    accent_right.line.color.rgb = COLORS["border"]
    at3 = s.shapes.add_textbox(Inches(11.45), Inches(4.27), Inches(1.0), Inches(0.5)).text_frame
    ap3 = at3.paragraphs[0]
    ap3.alignment = PP_ALIGN.CENTER
    ar3 = ap3.add_run()
    ar3.text = "RBAC\nSecure"
    ar3.font.size = Pt(11)
    ar3.font.bold = True
    ar3.font.color.rgb = COLORS["navy"]

    # Submitted by panel
    submit_panel = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.45), Inches(11.95), Inches(1.35))
    submit_panel.fill.solid()
    submit_panel.fill.fore_color.rgb = COLORS["white"]
    submit_panel.line.color.rgb = COLORS["border"]
    sh = s.shapes.add_textbox(Inches(1.0), Inches(5.65), Inches(2.0), Inches(0.35)).text_frame.paragraphs[0]
    sr = sh.add_run()
    sr.text = "Submitted by:"
    sr.font.size = Pt(14)
    sr.font.bold = True
    sr.font.color.rgb = COLORS["blue"]
    names = s.shapes.add_textbox(Inches(1.0), Inches(5.98), Inches(11.4), Inches(0.7)).text_frame
    np = names.paragraphs[0]
    np.alignment = PP_ALIGN.LEFT
    nr = np.add_run()
    nr.text = "Avanish Dubey | Ramlal Singh | Priyanshu Bhardwaj | Akash Kumar"
    nr.font.size = Pt(16)
    nr.font.bold = True
    nr.font.color.rgb = COLORS["navy"]

    # 2. Problem
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, COLORS["white"])
    add_title(s, "Problem Statement", "Why hospital teams need a policy intelligence agent")
    add_bullets(
        s,
        [
            "Hospital SOPs, circulars, billing policies, and claim rules are scattered across PDFs and folders.",
            "Teams lose time in manual lookup during admission, discharge, and claim processing.",
            "Incorrect or incomplete policy interpretation leads to rejection, delays, and compliance risk.",
            "Staff needs multilingual support (text + voice) for real-time operations.",
        ],
    )

    # 3. Solution
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, COLORS["white"])
    add_title(s, "Our Solution", "Policy Sarthi as a grounded, multilingual hospital operations copilot")
    add_card(s, Inches(0.75), Inches(1.65), Inches(3.9), Inches(2.25), "Ingest", "Admin uploads structured and unstructured policy data.", COLORS["soft_blue"])
    add_card(s, Inches(4.95), Inches(1.65), Inches(3.9), Inches(2.25), "Understand", "RAG + language + workflow context build grounded responses.", COLORS["soft_green"])
    add_card(s, Inches(9.15), Inches(1.65), Inches(3.35), Inches(2.25), "Assist", "Text + voice answers with source-backed guidance.", COLORS["soft_orange"])
    add_bullets(
        s,
        [
            "Role-based controls: only admin can upload policies.",
            "Feedback loop: thumbs up/down with correction capture.",
            "Designed for staff productivity, claim quality, and audit readiness.",
        ],
        top=4.35,
        height=2.2,
        size=16,
    )

    # 4. APIs
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, COLORS["white"])
    add_title(s, "Sarvam API Usage", "Current integration footprint")
    api_count = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.45), Inches(2.4), Inches(1.0))
    api_count.fill.solid()
    api_count.fill.fore_color.rgb = COLORS["blue"]
    api_count.line.fill.background()
    txt = s.shapes.add_textbox(Inches(1.05), Inches(1.68), Inches(1.9), Inches(0.5)).text_frame.paragraphs[0]
    rr = txt.add_run()
    rr.text = "5 APIs"
    rr.font.size = Pt(28)
    rr.font.bold = True
    rr.font.color.rgb = COLORS["white"]
    add_card(s, Inches(3.45), Inches(1.45), Inches(2.9), Inches(1.65), "1) Chat Completion", "Answer generation over retrieved context.", COLORS["soft_blue"])
    add_card(s, Inches(6.55), Inches(1.45), Inches(2.9), Inches(1.65), "2) Translation", "Non-English query/response handling.", COLORS["soft_green"])
    add_card(s, Inches(9.65), Inches(1.45), Inches(2.9), Inches(1.65), "3) Speech-to-Text", "Voice query transcription.", COLORS["soft_orange"])
    add_card(s, Inches(3.45), Inches(3.35), Inches(2.9), Inches(1.65), "4) Text-to-Speech", "Voice response playback.", COLORS["soft_orange"])
    add_card(s, Inches(6.55), Inches(3.35), Inches(2.9), Inches(1.65), "5) Document Intelligence", "OCR/text extraction from uploads.", COLORS["soft_blue"])
    add_bullets(
        s,
        [
            "Sarvam APIs are orchestrated in backend service layer (`sarvam_client.py`).",
            "Fallback paths exist when external API is unavailable.",
        ],
        top=5.35,
        height=1.2,
        size=14,
    )

    # 5. Workflow
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, COLORS["white"])
    add_title(s, "End-to-End Workflow", "From upload to grounded multilingual answer")
    add_flow_box(s, 0.6, 2.2, 2.0, 1.0, "Admin Upload\n(PDF/CSV/JSON)", COLORS["soft_blue"])
    add_flow_box(s, 2.9, 2.2, 2.0, 1.0, "Extraction +\nChunking/Indexing", COLORS["soft_green"])
    add_flow_box(s, 5.2, 2.2, 2.0, 1.0, "User Query\n(Text/Voice)", COLORS["soft_orange"])
    add_flow_box(s, 7.5, 2.2, 2.0, 1.0, "Retrieve +\nGround Answer", COLORS["soft_blue"])
    add_flow_box(s, 9.8, 2.2, 2.9, 1.0, "Response + Citation\n+ Optional Voice", COLORS["soft_green"])
    connect(s, 2.6, 2.7, 2.9, 2.7)
    connect(s, 4.9, 2.7, 5.2, 2.7)
    connect(s, 7.2, 2.7, 7.5, 2.7)
    connect(s, 9.5, 2.7, 9.8, 2.7)
    add_bullets(
        s,
        [
            "Structured data (CSV/JSON) is flattened and indexed as searchable rows.",
            "Unstructured data (SOPs/PDFs) is chunked for retrieval evidence.",
            "Feedback signal updates future retrieval ranking.",
        ],
        top=4.15,
        height=1.8,
        size=14,
    )

    # 6. Architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, COLORS["white"])
    add_title(s, "System Architecture", "Production-aligned layered design")
    add_flow_box(s, 1.2, 1.6, 10.9, 0.78, "Presentation Layer: Streamlit UI (chat, voice input, admin upload, feedback)", COLORS["soft_blue"])
    add_flow_box(s, 1.2, 2.6, 10.9, 0.78, "Application Layer: Flask APIs, AuthN/AuthZ, Routing, Validation", COLORS["soft_green"])
    add_flow_box(s, 1.2, 3.6, 10.9, 0.78, "AI Layer: RAG Orchestrator, Sarvam Client, Language Handling, Prompt Composer", COLORS["soft_orange"])
    add_flow_box(s, 1.2, 4.6, 10.9, 0.78, "Data Layer: SQLite (users/docs/chunks/feedback/structured rows) + File Storage", COLORS["soft_blue"])
    add_flow_box(s, 1.2, 5.6, 10.9, 0.78, "External Layer: Sarvam APIs + future vector DB + observability", COLORS["soft_green"])
    for y in [2.38, 3.38, 4.38, 5.38]:
        connect(s, 6.65, y, 6.65, y + 0.2)

    # 7. Data Strategy
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, COLORS["white"])
    add_title(s, "Data Strategy: Structured + Unstructured", "Unified retrieval to answer mixed policy questions")
    add_card(s, Inches(0.8), Inches(1.7), Inches(5.9), Inches(3.5), "Unstructured Pipeline", "Inputs: PDF, DOCX, text policies\nOCR/Extraction using Sarvam Document Intelligence\nChunking + metadata tagging\nRetrieved as sourceChunks for grounded answers", COLORS["soft_blue"])
    add_card(s, Inches(6.75), Inches(1.7), Inches(5.8), Inches(3.5), "Structured Pipeline", "Inputs: CSV/JSON datasets\nRow flattening and normalized search_text index\nStored as structured_records\nMatched rows added as evidence sections in answer context", COLORS["soft_green"])

    # 8. Security + Governance
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, COLORS["white"])
    add_title(s, "Security, Governance, and Reliability", "What the jury may ask about trust and controls")
    add_bullets(
        s,
        [
            "Role-based access control: admin upload restricted (`/api/documents/upload`).",
            "Token-based authentication for protected APIs.",
            "Feedback table + query logs create audit trail of behavior and quality signals.",
            "Grounded-response pattern minimizes hallucination via source-backed retrieval.",
            "Fallback mechanisms ensure graceful degradation when external AI services fail.",
        ],
    )

    # 9. Deployment
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, COLORS["white"])
    add_title(s, "Deployment Plan", "Where and how Policy Sarthi will run")
    add_card(s, Inches(0.8), Inches(1.7), Inches(3.8), Inches(2.4), "MVP Deployment", "Single VM/Container\nFrontend + Backend behind reverse proxy\nManaged secrets for Sarvam key", COLORS["soft_blue"])
    add_card(s, Inches(4.95), Inches(1.7), Inches(3.8), Inches(2.4), "Preferred Cloud", "Azure / AWS / GCP\nApp service + managed DB + object storage\nCI/CD via GitHub Actions", COLORS["soft_green"])
    add_card(s, Inches(9.1), Inches(1.7), Inches(3.4), Inches(2.4), "Hospital Integration", "VPN/private network\nSSO with hospital IAM\nAudit export to SIEM", COLORS["soft_orange"])
    add_bullets(
        s,
        [
            "Target environments: staging (UAT) -> production with change approvals.",
            "Secrets, logs, and backups managed through cloud-native services.",
        ],
        top=4.45,
        height=1.3,
        size=14,
    )

    # 10. Scale
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, COLORS["white"])
    add_title(s, "Scalability Roadmap", "How we scale from pilot to enterprise")
    add_card(s, Inches(0.8), Inches(1.7), Inches(3.9), Inches(3.6), "Phase 1 (Now)", "SQLite + local storage\nSingle service deployment\nHospital-level pilot", COLORS["soft_blue"])
    add_card(s, Inches(4.95), Inches(1.7), Inches(3.9), Inches(3.6), "Phase 2", "Postgres + Redis caching\nVector DB (Milvus/Pinecone)\nAsync workers for ingestion\nMulti-hospital tenancy", COLORS["soft_green"])
    add_card(s, Inches(9.1), Inches(1.7), Inches(3.4), Inches(3.6), "Phase 3", "Horizontal autoscaling\nRegional failover\nModel routing + cost controls\nSLA/SLO observability", COLORS["soft_orange"])

    # 11. KPIs
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, COLORS["white"])
    add_title(s, "Success Metrics", "What we will measure post-deployment")
    add_card(s, Inches(0.8), Inches(1.7), Inches(3.9), Inches(2.4), "Operational KPIs", "Avg response latency\nPolicy lookup time saved\nDaily active staff users", COLORS["soft_blue"])
    add_card(s, Inches(4.95), Inches(1.7), Inches(3.9), Inches(2.4), "Quality KPIs", "Grounded answer rate\nThumbs-up ratio\nClaim rejection reduction", COLORS["soft_green"])
    add_card(s, Inches(9.1), Inches(1.7), Inches(3.4), Inches(2.4), "Reliability KPIs", "API uptime\nError rate\nFallback usage rate", COLORS["soft_orange"])
    add_bullets(
        s,
        [
            "Feedback loop is already implemented and directly influences retrieval ranking.",
            "Next: automated evaluation set for regression testing before releases.",
        ],
        top=4.45,
        height=1.2,
        size=14,
    )

    # 12. Jury Q&A prep
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, COLORS["white"])
    add_title(s, "Jury Questions We Are Ready For", "Technical depth and product-readiness")
    add_bullets(
        s,
        [
            "How do you prevent hallucination? -> Grounded retrieval + source evidence + no-info fallback.",
            "How is multilingual handled? -> Sarvam translation + language detection + voice STT/TTS.",
            "How do you secure hospital data? -> Auth, RBAC, controlled upload, audit logs.",
            "How will this scale beyond one hospital? -> DB upgrade, vector DB, async ingestion, multi-tenant architecture.",
            "What is the ROI? -> Faster claim processing, lower rejection, reduced policy lookup time.",
        ],
        size=16,
    )

    # 13. Closing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, COLORS["bg"])
    add_title(s, "Thank You", "Policy Sarthi AI: From static documents to real-time hospital intelligence")
    chip = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.2))
    chip.fill.solid()
    chip.fill.fore_color.rgb = COLORS["navy"]
    chip.line.fill.background()
    p = s.shapes.add_textbox(Inches(1.4), Inches(2.6), Inches(10.5), Inches(0.5)).text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Demo Flow: Upload Policy -> Ask in Any Language -> Get Grounded Answer + Voice + Feedback"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]
    add_footer(s, "Prepared for jury evaluation: architecture, APIs, governance, deployment, and scale.")

    prs.save(OUTPUT_FILE)


if __name__ == "__main__":
    build()
