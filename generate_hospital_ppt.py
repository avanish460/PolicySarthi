from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_FILE = "AI_Powered_Hospital_Ayushman_Assistant.pptx"


COLORS = {
    "navy": RGBColor(20, 61, 89),
    "teal": RGBColor(31, 138, 112),
    "light": RGBColor(244, 246, 248),
    "white": RGBColor(255, 255, 255),
    "orange": RGBColor(242, 140, 40),
    "gray": RGBColor(92, 101, 112),
    "dark": RGBColor(38, 50, 56),
    "border": RGBColor(208, 214, 220),
    "soft_blue": RGBColor(228, 238, 247),
    "soft_green": RGBColor(227, 242, 237),
    "soft_orange": RGBColor(252, 239, 225),
}


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.2), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = COLORS["navy"]
    p.alignment = PP_ALIGN.LEFT

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(2.4), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["teal"]
    bar.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.2), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.size = Pt(11)
        r.font.color.rgb = COLORS["gray"]


def add_bullet_slide(slide, title, lead, bullets):
    add_title(slide, title)

    lead_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(11.6), Inches(0.9))
    tf = lead_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = lead
    r.font.size = Pt(16)
    r.font.color.rgb = COLORS["dark"]

    body = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.1), Inches(4.7))
    tf = body.text_frame
    tf.word_wrap = True
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS["navy"]
        p.bullet = True


def add_footer(slide, text):
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.0), Inches(0.3))
    p = footer.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = COLORS["gray"]


def add_card(slide, left, top, width, height, title, body, fill):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = COLORS["border"]

    title_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(15)
    r.font.color.rgb = COLORS["navy"]

    body_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.45), width - Inches(0.3), height - Inches(0.55))
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = body
    r.font.size = Pt(11)
    r.font.color.rgb = COLORS["dark"]


def add_layer(slide, left, top, width, height, title, fill, items):
    layer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    layer.fill.solid()
    layer.fill.fore_color.rgb = fill
    layer.line.color.rgb = COLORS["border"]

    title_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.28))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(15)
    r.font.color.rgb = COLORS["navy"]
    p.alignment = PP_ALIGN.CENTER

    item_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.42), width - Inches(0.3), height - Inches(0.55))
    tf = item_box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS["dark"]
        p.alignment = PP_ALIGN.CENTER


def add_vertical_arrow(slide, x_center, top, bottom):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x_center, top, x_center, bottom)
    line.line.color.rgb = COLORS["teal"]
    line.line.width = Pt(2)
    line.line.end_arrowhead = True


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLORS["navy"]
    banner.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(11.8), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "PolicySaarthi"
    r.font.bold = True
    r.font.size = Pt(27)
    r.font.color.rgb = COLORS["navy"]

    sub = slide.shapes.add_textbox(Inches(0.65), Inches(2.75), Inches(11.2), Inches(0.6))
    p = sub.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Enabling faster policy access, accurate insurance guidance, and multilingual support for hospital staff"
    r.font.size = Pt(18)
    r.font.color.rgb = COLORS["gray"]

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(4.0), Inches(3.8), Inches(0.65))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["teal"]
    accent.line.fill.background()
    p = slide.shapes.add_textbox(Inches(0.95), Inches(4.16), Inches(3.2), Inches(0.3)).text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Built for document-heavy hospital workflows"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]

    team = slide.shapes.add_textbox(Inches(0.7), Inches(6.4), Inches(5.5), Inches(0.5))
    p = team.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Team Name | Member Names | Hackathon Name"
    r.font.size = Pt(13)
    r.font.color.rgb = COLORS["gray"]

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_bullet_slide(
        slide,
        "Problem We Are Solving",
        "Hospitals operate with thousands of SOPs, circulars, discharge policies, and insurance guidelines spread across PDFs, scans, and folders. During live patient interactions, staff need answers quickly but the information is hard to access.",
        [
            "Manual document search slows down admissions, discharge, and claim workflows",
            "Ayushman Bharat claim processing is difficult because eligibility, required documents, and pre-authorization steps are not easy to locate",
            "Staff often depend on senior colleagues instead of a reliable, searchable knowledge system",
            "English-only documents create communication barriers for multilingual staff and patients",
        ],
    )
    add_footer(slide, "Bottom line: hospital teams need an AI assistant that turns scattered documents into operational guidance.")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_title(slide, "Why This Problem Matters")
    add_card(slide, Inches(0.8), Inches(1.5), Inches(2.8), Inches(1.55), "Operational Delay", "Staff lose valuable time searching for the latest process and the right document section.", COLORS["soft_blue"])
    add_card(slide, Inches(3.95), Inches(1.5), Inches(2.8), Inches(1.55), "Claim Rejections", "Missing documents or incomplete workflow understanding causes rework and rejected claims.", COLORS["soft_orange"])
    add_card(slide, Inches(7.1), Inches(1.5), Inches(2.8), Inches(1.55), "Language Barrier", "Many hospital staff and patients are more comfortable in Hindi or a regional language.", COLORS["soft_green"])
    add_card(slide, Inches(10.25), Inches(1.5), Inches(2.25), Inches(1.55), "Compliance Risk", "Using outdated SOPs or wrong steps affects governance and service quality.", COLORS["soft_blue"])
    quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(4.0), Inches(11.5), Inches(1.25))
    quote.fill.solid()
    quote.fill.fore_color.rgb = COLORS["navy"]
    quote.line.fill.background()
    qtf = slide.shapes.add_textbox(Inches(1.25), Inches(4.35), Inches(10.8), Inches(0.5)).text_frame
    p = qtf.paragraphs[0]
    r = p.add_run()
    r.text = "A small document-access problem creates a big business, compliance, and patient-service problem."
    r.font.size = Pt(21)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_bullet_slide(
        slide,
        "Our Solution",
        "We propose an AI-powered multilingual document and insurance assistant for hospitals.",
        [
            "Ingest hospital SOPs, claim guidelines, discharge policies, and Ayushman workflow documents",
            "Use OCR to extract text from scanned files and make them searchable",
            "Answer staff questions through text or voice in English, Hindi, and regional languages",
            "Provide grounded answers with source-backed citations from uploaded documents",
            "Guide insurance desks through Ayushman claim steps, document checklists, and process clarifications",
        ],
    )
    add_footer(slide, "From document chaos to guided hospital decisions.")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_title(slide, "Key Features")
    features = [
        ("Document Upload & OCR", "Upload PDFs, scans, forms, circulars, and SOPs; extract searchable text."),
        ("Semantic Search", "Retrieve the right clause or policy section instantly."),
        ("Multilingual Chat", "Ask questions in English, Hindi, or another supported language."),
        ("Voice Assistance", "Enable speech-to-text and text-to-speech for fast hands-free use."),
        ("Ayushman Claim Guidance", "Get step-by-step assistance for eligibility, documents, and claim flow."),
        ("Source Citation", "Every answer links back to the originating hospital document."),
    ]
    positions = [
        (Inches(0.7), Inches(1.5)),
        (Inches(4.45), Inches(1.5)),
        (Inches(8.2), Inches(1.5)),
        (Inches(0.7), Inches(3.65)),
        (Inches(4.45), Inches(3.65)),
        (Inches(8.2), Inches(3.65)),
    ]
    fills = [COLORS["soft_blue"], COLORS["soft_green"], COLORS["soft_orange"], COLORS["soft_green"], COLORS["soft_blue"], COLORS["soft_orange"]]
    for (title, body), (left, top), fill in zip(features, positions, fills):
        add_card(slide, left, top, Inches(3.4), Inches(1.75), title, body, fill)
    add_footer(slide, "Future-ready for version comparison, analytics, and audit trails.")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_title(slide, "Ayushman Claim Workflow Use Case")
    use_case = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(5.2), Inches(4.2))
    use_case.fill.solid()
    use_case.fill.fore_color.rgb = COLORS["soft_blue"]
    use_case.line.color.rgb = COLORS["border"]
    ub = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(4.7), Inches(3.5))
    tf = ub.text_frame
    tf.word_wrap = True
    lines = [
        "Patient arrives with an Ayushman Card",
        "Insurance desk must verify eligibility and required documents",
        "Team checks treatment/package guidance and pre-authorization needs",
        "System returns step-by-step instructions from hospital SOPs",
        "Answer is shown in English or Hindi with document evidence",
    ]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.bullet = True
        p.font.size = Pt(18 if idx == 0 else 16)
        p.font.color.rgb = COLORS["navy"]
        if idx == 0:
            p.font.bold = True

    query_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.35), Inches(1.7), Inches(5.95), Inches(1.0))
    query_box.fill.solid()
    query_box.fill.fore_color.rgb = COLORS["teal"]
    query_box.line.fill.background()
    tf = slide.shapes.add_textbox(Inches(6.65), Inches(2.0), Inches(5.35), Inches(0.5)).text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = '"What documents are required for Ayushman claim submission?"'
    r.font.size = Pt(17)
    r.font.color.rgb = COLORS["white"]
    r.font.bold = True

    response = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.35), Inches(3.0), Inches(5.95), Inches(2.35))
    response.fill.solid()
    response.fill.fore_color.rgb = COLORS["soft_green"]
    response.line.color.rgb = COLORS["border"]
    tf = slide.shapes.add_textbox(Inches(6.65), Inches(3.3), Inches(5.35), Inches(1.8)).text_frame
    tf.word_wrap = True
    for idx, text in enumerate(
        [
            "Required document checklist",
            "Pre-authorization and submission steps",
            "Relevant hospital SOP reference",
            "Hindi / English multilingual response",
        ]
    ):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.bullet = True
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["dark"]

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_title(slide, "System Architecture", "Layered architecture for multilingual document intelligence and claim guidance")
    layer_specs = [
        ("USER INTERACTION", COLORS["soft_blue"], ["Hospital Staff", "Admin", "Insurance Desk", "Web / Mobile / Chat / Voice"]),
        ("APPLICATION & ACCESS LAYER", COLORS["light"], ["API Gateway", "Authentication", "Role-Based Access", "Audit Logging"]),
        ("AI ORCHESTRATION LAYER", COLORS["soft_green"], ["Query Router", "Language Detection", "Prompt Builder", "RAG Workflow"]),
        ("DOCUMENT INTELLIGENCE PIPELINE", COLORS["soft_orange"], ["Upload Docs", "OCR Extraction", "Chunking", "Metadata Tagging"]),
        ("DATA LAYER", COLORS["light"], ["Vector Database", "Document Store", "App DB", "Users / Logs / Feedback"]),
    ]
    top = 1.35
    height = 0.92
    left = Inches(1.1)
    width = Inches(11.1)
    centers = []
    for idx, (title, fill, items) in enumerate(layer_specs):
        current_top = Inches(top + idx * 1.05)
        add_layer(slide, left, current_top, width, Inches(height), title, fill, items)
        centers.append(current_top + Inches(height))
    for idx in range(len(centers) - 1):
        add_vertical_arrow(slide, Inches(6.65), centers[idx], centers[idx + 1] - Inches(0.12))

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_title(slide, "How It Works")
    steps = [
        "Hospital uploads SOPs, claim guidelines, and Ayushman documents",
        "OCR extracts text from scanned files and images",
        "Documents are chunked, tagged, and indexed in the vector database",
        "Staff asks a question through chat or voice",
        "System retrieves relevant sections and generates a grounded answer",
        "Response is returned in English or Hindi with source citation",
    ]
    x = 0.55
    for idx, step in enumerate(steps, start=1):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3), Inches(1.92), Inches(1.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = [COLORS["soft_blue"], COLORS["soft_green"], COLORS["soft_orange"]][(idx - 1) % 3]
        shape.line.color.rgb = COLORS["border"]
        num = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.7), Inches(1.8), Inches(0.5), Inches(0.5))
        num.fill.solid()
        num.fill.fore_color.rgb = COLORS["navy"]
        num.line.fill.background()
        p = slide.shapes.add_textbox(Inches(x + 0.84), Inches(1.92), Inches(0.25), Inches(0.2)).text_frame.paragraphs[0]
        r = p.add_run()
        r.text = str(idx)
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = COLORS["white"]
        tf = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.55), Inches(1.62), Inches(0.9)).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = step
        r.font.size = Pt(11)
        r.font.color.rgb = COLORS["dark"]
        if idx < len(steps):
            add_vertical_arrow(slide, Inches(x + 1.95), Inches(3.0), Inches(3.0))
        x += 2.1
    add_footer(slide, 'Example query: "Patient ke Ayushman claim ke liye kaunse documents chahiye?"')

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_title(slide, "Technology Stack & API Mapping")
    add_card(slide, Inches(0.7), Inches(1.6), Inches(2.8), Inches(2.2), "Frontend", "React or Next.js web app\nChat UI\nVoice input button\nAdmin upload dashboard", COLORS["soft_blue"])
    add_card(slide, Inches(3.85), Inches(1.6), Inches(2.8), Inches(2.2), "Backend", "Node.js or Python service\nAuth and role management\nRAG pipeline\nResponse orchestration", COLORS["soft_green"])
    add_card(slide, Inches(7.0), Inches(1.6), Inches(2.8), Inches(2.2), "Data Layer", "Vector database\nSQL or NoSQL metadata store\nBlob storage for documents", COLORS["soft_orange"])
    add_card(slide, Inches(10.15), Inches(1.6), Inches(2.45), Inches(2.2), "Sarvam APIs", "Chat Completion\nTranslation\nSpeech-to-Text\nText-to-Speech\nDocument Intelligence", COLORS["soft_blue"])
    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(4.45), Inches(10.9), Inches(1.1))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLORS["navy"]
    banner.line.fill.background()
    p = slide.shapes.add_textbox(Inches(1.5), Inches(4.78), Inches(10.3), Inches(0.5)).text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Strong fit for Sarvam: document extraction + multilingual intelligence + voice + grounded chat"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLORS["white"])
    add_title(slide, "Impact & Future Scope")
    add_card(slide, Inches(0.8), Inches(1.7), Inches(5.65), Inches(3.3), "Impact", "Faster claim processing\nFewer documentation mistakes\nReduced claim rejection\nImproved compliance\nBetter multilingual support\nFaster staff onboarding", COLORS["soft_green"])
    add_card(slide, Inches(6.85), Inches(1.7), Inches(5.65), Inches(3.3), "Future Scope", "Claim rejection prediction\nPolicy version comparison\nPatient-facing assistant\nAnalytics dashboard\nHospital management system integration\nExpansion to other insurance workflows", COLORS["soft_blue"])
    closing = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(5.55), Inches(11.2), Inches(0.95))
    closing.fill.solid()
    closing.fill.fore_color.rgb = COLORS["teal"]
    closing.line.fill.background()
    p = slide.shapes.add_textbox(Inches(1.45), Inches(5.85), Inches(10.4), Inches(0.3)).text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Transforming static hospital documents into a multilingual AI assistant for real-time operational decisions."
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]

    prs.save(OUTPUT_FILE)


if __name__ == "__main__":
    build_presentation()
